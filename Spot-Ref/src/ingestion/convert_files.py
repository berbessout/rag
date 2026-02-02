import pytesseract
from pdf2image import convert_from_path
import cv2
import numpy as np
from pathlib import Path
from pptx import Presentation


def convert_pdf_to_txt(filename: str, input_path: Path) -> str | None:
    """
    Convert a PDF file to a text file using OCR.
    
    Args:
        filename: Name of the PDF file
        input_path: Directory containing the PDF file
        output_path: Directory where to save the text file
    """

    pdf_path = input_path / filename

    if not pdf_path.exists():
        print(f"❌ Fichier PDF non trouvé : {pdf_path}")
        return
    print(f"🌀 Converting : {filename}")

    try:
        # Convert the PDF to images
        images = convert_from_path(str(pdf_path))
        text = ""
        for i, image in enumerate(images):
            # Convert the PIL image to OpenCV format
            cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)

            # Extract all the text from the page
            custom_config = r'--oem 3 --psm 6 -l fra'
            text += pytesseract.image_to_string(cv_image, config=custom_config)
            text += "\n"
        return text

    except Exception as e:
        print(f"❌ Error converting {filename} : {e}")
        return None
    
def convert_pptx_to_txt(filename: str, input_path: Path) -> str | None:
    """
    Convert a PPTX file to a text file using python-pptx.
    
    Args:
        filename: Name of the PPTX file
        input_path: Directory containing the PPTX file
        output_dir: Directory where to save the text file
        
    Returns:
        Path to the created text file, or None if conversion failed

    pre-condition: filename is a valid PPTX file
    post-condition: the text is either returned or None if conversion failed
    """

    input_file = input_path / filename

    if not input_file.exists():
        print(f"❌ PPTX file not found : {input_file}")
        return None

    print(f"🌀 Converting : {filename}")

    try:
        prs = Presentation(str(input_file))
        text_content = []

        for slide_number, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n=== Slide {slide_number} ===\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        text_content = "\n".join(text_content)
        return text_content

    except Exception as e:
        print(f"❌ Error converting {filename} : {e}")
        return None

def convert_files(filename: str, input_path: Path) -> str | None:
    """
    retrieve a file text from a file and return it

    Args:
        filename: Name of the file
        input_path: Directory containing the file

    return:
        the text of the file or None if conversion failed

    pre-condition: filename is a valid file
    post-condition: the text is either returned or None if conversion failed
    """
    file = input_path / filename
    if file.suffix in ['.pptx', '.ppt']:
        raw_text = convert_pptx_to_txt(filename, input_path)
    elif file.suffix in ['.pdf']:
        raw_text = convert_pdf_to_txt(filename, input_path)
    else:
        print(f"📄 File format not supported : {file.suffix}")
        return None
    return raw_text